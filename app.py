import os
import backoff 
import streamlit as st
import openai
import pptx

#API Keys
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
if not OPENAI_API_KEY:
    raise ValueError("Missing OPENAI_API_KEY environment variable.")
openai.api_key = OPENAI_API_KEY

@backoff.on_exception(backoff.expo, openai.error.RateLimitError)
def get_gpt4_response(system_prompt, user_prompt):
    """Returns a davinci answer from OpenAI"""
    gpt_response = openai.ChatCompletion.create(
        model="gpt-4",
        temperature=0.00,
        max_tokens=2000,
        messages = [{"role": "system", "content": system_prompt}, {"role": "user", "content" : user_prompt}]
    )
    return gpt_response['choices'][0]['message']['content'].strip()

def extract_pptx_text(pptx_file):
    prs = pptx.Presentation(pptx_file)
    text = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text.append(paragraph.text) 
    return '\n'.join(text)

st.title("Convert a PowerPoint into a draft speech!")
st.write("Upload a .ppt file and get a draft speech in return that you can customise")
st.warning("Please note that the contents of the file should be less than 4000 words in length. Once you upload the file please wait for the completion.")
pptx_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"])
if pptx_file is not None:
    document_text = extract_pptx_text(pptx_file)
    if len(document_text.split()) > 4000:
        st.error("The uploaded file is too long. Please upload a file with less than 4000 words.")
    else:
        system_prompt = "You are an expert speech writer"
        user_prompt = f"Use the following information from the provided PPT to create a captivating speech: {document_text}"
        gpt4_response = get_gpt4_response(system_prompt, user_prompt)
        st.write("The uploaded file has been successfully processed. Your draft speech is below:")
        css = """
            <style>
                code {
                    white-space: pre-wrap !important;
                    primaryColor="#d33682";
                    backgroundColor="#000000";
                    secondaryBackgroundColor="#2c738a";
                    textColor="#fafafa";
                    font="serif";
                    padding: 10px;
                    border-radius: 5px;
                    font-size: 14px;
                }
            </style>
        """
        st.markdown(css, unsafe_allow_html=True)
        st.code(gpt4_response, language="markdown")
        
